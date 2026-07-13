# -*- coding: utf-8 -*-
# Build the MES Demo Evaluation deck: title -> per-section divider -> one slide per Section-A point.
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
sys.path.insert(0, os.path.dirname(__file__))
import shots

IMG = os.path.join(os.path.dirname(__file__), "img")
OUT = os.path.join(os.path.dirname(__file__), "MES_Demo_Evaluation_TrackA.pptx")

PRIMARY = RGBColor(0x1E, 0x8F, 0xBF); ACCENT = RGBColor(0x25, 0xA9, 0xE0)
NAVY = RGBColor(0x14, 0x1B, 0x2E); DARK = RGBColor(0x1A, 0x22, 0x3A)
INK = RGBColor(0x22, 0x2A, 0x35); MUT = RGBColor(0x5A, 0x66, 0x74); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32); AMBER = RGBColor(0xC4, 0x6A, 0x00); FRAME = RGBColor(0xD5, 0xDC, 0xE3)
PANEL = RGBColor(0xF4, 0xF7, 0xFA); PBLUE = RGBColor(0xEA, 0xF4, 0xFA)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bgcolor=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bgcolor
    return s

def rect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line: sp.line.color.rgb = line; sp.line.width = Pt(1)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.0):
    # runs: list of paragraphs; each paragraph = list of (str, size, color, bold)
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in (tf.margin_left, ): pass
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line_sp
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt; f = r.font
            f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = "Calibri"
    return tb

def pic_fit(s, path, l, t, w, h):
    try:
        im = Image.open(path); iw, ih = im.size
    except Exception:
        rect(s, l, t, w, h, PANEL, FRAME); return
    ar = iw / ih; bar = w / h
    if ar > bar: nw, nh = w, w / ar
    else: nh, nw = h, h * ar
    x, y = l + (w - nw) / 2, t + (h - nh) / 2
    rect(s, x - 0.03, y - 0.03, nw + 0.06, nh + 0.06, FRAME)   # thin frame
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(nw), Inches(nh))

def pill(s, l, t, w, h, label, color):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False
    tf = sp.text_frame; tf.margin_top = tf.margin_bottom = Inches(0.01); tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = label
    r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    return sp

# ---------------- Title slide ----------------
s = slide(NAVY)
rect(s, 0, 0, SW, 0.28, ACCENT)
rect(s, 0, SH - 0.22, SW, 0.22, PRIMARY)
text(s, 0.9, 1.5, 11.5, 1.2, [[("Bluemingo MES ", 15, ACCENT, True), ("v2", 15, RGBColor(0xAF,0xC4,0xD6), False)]])
text(s, 0.9, 2.15, 11.6, 2.2, [
    [("Integrated Steel Works", 40, WHITE, True)],
    [("MES · Functional Demo Walkthrough", 24, RGBColor(0xCF,0xDD,0xEA), False)],
])
text(s, 0.9, 4.25, 11.6, 1.4, [
    [("Track A — Functional Fit Demonstration", 18, ACCENT, True)],
    [("Every functional point, mapped to the live MES screen, execution path and system response.", 14, RGBColor(0xB9,0xC8,0xD6), False)],
])
text(s, 0.9, 6.35, 11.6, 0.6, [[
    ("SMS-2 (KR → BOF → ARS → LHF → RH → CCM)   ·   HSM-2 Hot Strip Mill   ·   SAP / Level-2 / SYMS-CYMS integration",
     12, RGBColor(0x9F,0xB2,0xC4), False)]])
text(s, 0.9, 6.95, 11.6, 0.4, [[("Tracked thread:  a 5,000 MT HRC order (IS 2062 E250)  →  Heat H2601  →  slabs  →  Coil C2601",
     11.5, RGBColor(0x86,0x9A,0xAD), False)]])

# ---------------- Agenda / coverage slide ----------------
s = slide(WHITE)
rect(s, 0, 0, SW, 1.0, DARK)
text(s, 0.6, 0.18, 12, 0.7, [[("What this deck covers", 26, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
counts = {}
for it in shots.ITEMS:
    counts[it["grp"]] = counts.get(it["grp"], 0) + 1
labels = dict(shots.SECTIONS)
y = 1.45
rows = [(g, labels[g], counts.get(g, 0)) for g, _ in shots.SECTIONS]
colw = 6.1
for i, (g, lab, n) in enumerate(rows):
    col = i // 4; row = i % 4
    x = 0.6 + col * (colw + 0.5); yy = 1.5 + row * 1.35
    rect(s, x, yy, colw, 1.15, PANEL, FRAME)
    rect(s, x, yy, 0.11, 1.15, ACCENT)
    text(s, x + 0.28, yy + 0.14, colw - 0.5, 0.9, [
        [(lab, 15, INK, True)],
        [("%d evaluation points" % n, 12, MUT, False)],
    ], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.6, 6.9, 12, 0.4, [[
    ("50 points across A1–A7 (Rake Yard excluded).  Each point = one slide: operator action · execution path · expected system response · live screenshot.",
     12, MUT, False)]])

# ---------------- Section dividers + item slides ----------------
def section_divider(grp, label):
    s = slide(DARK)
    rect(s, 0, 0, 0.35, SH, ACCENT)
    text(s, 1.1, 2.5, 11, 1.0, [[(grp, 60, RGBColor(0x2E,0x5A,0x78), True)]])
    text(s, 1.15, 3.7, 11, 1.4, [[(label.split("·", 1)[1].strip() if "·" in label else label, 30, WHITE, True)]])
    subs = []
    for it in shots.ITEMS:
        if it["grp"] == grp and it["sec"] not in subs:
            subs.append(it["sec"])
    text(s, 1.18, 5.0, 11, 1.8, [[(" •  " + ss.split("·", 1)[1].strip(), 14, RGBColor(0xB9,0xC8,0xD6), False)] for ss in subs])

def item_slide(it):
    s = slide(WHITE)
    # header band
    rect(s, 0, 0, SW, 1.28, DARK)
    rect(s, 0, 1.28, SW, 0.06, ACCENT)
    prc = GREEN if it["pr"] == "P1" else AMBER
    pill(s, 0.55, 0.24, 0.95, 0.34, it["id"], PRIMARY)
    pill(s, 1.60, 0.24, 0.62, 0.34, it["pr"], prc)
    text(s, 0.55, 0.66, 12.2, 0.6, [[(it["title"], 17.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    # left panel: sub-section, execution path, expected result
    lx, lw = 0.45, 4.15; ty = 1.62
    text(s, lx, ty, lw, 0.35, [[(it["sec"].split("·", 1)[1].strip() if "·" in it["sec"] else it["sec"], 12, ACCENT, True)]])
    text(s, lx, ty + 0.42, lw, 0.3, [[("EXECUTION PATH", 11, MUT, True)]])
    steps = []
    for i, stp in enumerate(it["path"], 1):
        steps.append([("%d. " % i, 12.5, ACCENT, True), (stp, 12.5, INK, False)])
    text(s, lx, ty + 0.76, lw, 2.6, steps, sp_after=6, line_sp=1.02)
    # expected result box
    ry = ty + 0.80 + 0.52 * len(it["path"]) + 0.25
    ry = min(ry, 5.2)
    boxh = 6.95 - ry
    rect(s, lx, ry, lw, boxh, PBLUE, RGBColor(0xC5,0xE2,0xF0))
    rect(s, lx, ry, 0.10, boxh, ACCENT)
    text(s, lx + 0.24, ry + 0.12, lw - 0.4, boxh - 0.2, [
        [("EXPECTED SYSTEM RESPONSE", 10.5, PRIMARY, True)],
        [(it["resp"], 12.5, INK, False)],
    ], sp_after=4, line_sp=1.03)
    # screenshot (right)
    pic_fit(s, os.path.join(IMG, it["id"] + ".png"), 4.85, 1.62, 8.05, 5.35)
    # footer
    text(s, 4.85, 7.02, 8.05, 0.3, [[("Bluemingo MES — Track A Functional Fit   ·   " + it["sec"].split("–")[0].strip(),
        10, MUT, False)]], align=PP_ALIGN.LEFT)

seen = set()
for it in shots.ITEMS:
    if it["grp"] not in seen:
        seen.add(it["grp"]); section_divider(it["grp"], dict(shots.SECTIONS)[it["grp"]])
    item_slide(it)

# ---------------- Closing ----------------
s = slide(NAVY)
rect(s, 0, 0, SW, 0.28, ACCENT)
text(s, 0.9, 2.2, 11.5, 1.0, [[("Functional fit — proven end to end", 34, WHITE, True)]])
bullets = [
    "Unified schema spans SMS-2 continuous casting and HSM-2 hot rolling — one data model, split & merge topologies.",
    "Live Level-2 (BOF / CCM / Mill, OPC-UA PDI/PDO) and SAP (SD / PP / QM / PM / MM, IDoc / BAPI) integration.",
    "Full traceability: heat genealogy, 14–20 char slab/coil IDs, batch-order allocation, non-repudiable audit trail.",
    "Governed quality: sampling, chemistry vs spec, out-of-spec auto-hold, downgrade + alternate-SO, AUD, FG recall.",
    "Reactive scheduling, delay capture, yard (SYMS/CYMS) allocation, and role-based access control.",
]
text(s, 0.95, 3.5, 11.4, 3.0, [[(" ✓  ", 15, ACCENT, True), (b, 15, RGBColor(0xDA,0xE6,0xF0), False)] for b in bullets], sp_after=8, line_sp=1.05)
text(s, 0.9, 6.7, 11.5, 0.5, [[("Bluemingo MES v2  ·  Track A functional fit demonstrated", 14, ACCENT, True)]])

prs.save(OUT)
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
